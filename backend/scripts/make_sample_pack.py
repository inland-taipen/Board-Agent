"""Generate a synthetic board pack for testing and demos.

Produces the four formats the BRD names, for a fictional listed company. The
content is deliberately constructed so a correct briefing has findings to make:

* Prior minutes carry actions that the current pack does NOT report back on -
  the reconciliation pass should mark those 'unclear'.
* The deck attributes the margin decline to one cause and the MIS to another.
* The risk register and the financials quantify the same exposure differently.
* One action from two meetings ago is still open and is restated in different
  words, which exercises action-item deduplication.

Nothing here is real. "Meridian Industries Limited" is fictional.

Usage:  python scripts/make_sample_pack.py ../sample_pack
"""

from __future__ import annotations

import sys
from pathlib import Path

import docx
import openpyxl
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PPt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

COMPANY = "Meridian Industries Limited"


# --- 1. Prior minutes (PDF) ---------------------------------------------------

MINUTES = [
    ("h1", "MINUTES OF THE 118TH MEETING OF THE BOARD OF DIRECTORS"),
    ("h2", f"{COMPANY} — held on 14 May 2026 at the registered office, Pune"),
    ("p", "Present: Mr R. Kulkarni (Chairman), Ms A. Desai (Independent Director), "
          "Mr S. Iyer (Independent Director), Ms P. Nambiar (Independent Director), "
          "Mr V. Meridian (Managing Director), Mr D. Shah (Chief Financial Officer). "
          "In attendance: Ms N. Rao, Company Secretary."),
    ("h2", "Item 4 — Financial results for the quarter ended 31 March 2026"),
    ("p", "The CFO presented the audited results. EBITDA margin for the quarter stood at "
          "17.8 per cent against 19.1 per cent in the corresponding quarter. The Board "
          "noted the results and took them on record."),
    ("p", "The Board expressed concern at the growth in trade receivables beyond 180 days. "
          "It was RESOLVED that the CFO shall place before the Board, at the next meeting, "
          "an ageing analysis of receivables above 180 days together with a recovery plan "
          "for the ten largest overdue accounts."),
    ("h2", "Item 5 — Report of the Audit Committee"),
    ("p", "Ms Desai reported that the internal audit of the Nashik plant had identified a "
          "repeat observation on segregation of duties in the procurement function, first "
          "raised in the audit cycle ended March 2025."),
    ("p", "The Board DIRECTED that management complete the remediation of the procurement "
          "segregation-of-duties observation and report closure to the Audit Committee by "
          "30 September 2026."),
    ("h2", "Item 6 — Information technology and cyber security"),
    ("p", "The Managing Director informed the Board that the IT security roadmap prepared "
          "by the consultants was under review. It was AGREED that the IT security roadmap, "
          "together with an indicative capital outlay, be placed before the Board at the "
          "next meeting."),
    ("p", "The Board further DIRECTED that a report on the company's cyber insurance "
          "coverage and its adequacy be circulated to directors within thirty days."),
    ("h2", "Item 7 — Kalyani Auto Components acquisition"),
    ("p", "The Board considered the proposal to acquire a 74 per cent stake in Kalyani Auto "
          "Components Private Limited. The proposal was DEFERRED. The Managing Director was "
          "REQUESTED to obtain an independent valuation opinion and to revert to the Board "
          "with a revised proposal."),
    ("h2", "Item 8 — Succession planning"),
    ("p", "Ms Nambiar raised the absence of a documented succession plan for key managerial "
          "personnel. The Board ADVISED management to prepare a succession plan for the "
          "Managing Director and Chief Financial Officer positions and to place it before "
          "the Nomination and Remuneration Committee."),
    ("h2", "Item 9 — Any other business"),
    ("p", "The Chairman noted that the action on standardising the group's ESG disclosure "
          "framework, raised at the 116th meeting held on 12 February 2026, remained "
          "outstanding. Management UNDERTOOK to revert at the next meeting."),
    ("p", "There being no other business, the meeting concluded with a vote of thanks to "
          "the Chair."),
    ("pb", ""),
    ("h1", "MINUTES OF THE 117TH MEETING OF THE BOARD OF DIRECTORS"),
    ("h2", f"{COMPANY} — held on 12 February 2026"),
    ("h2", "Item 3 — Environmental, social and governance reporting"),
    ("p", "The Board considered the group's ESG disclosure practice. It was RESOLVED that "
          "management shall standardise the ESG disclosure framework across all subsidiaries "
          "and present a consolidated framework to the Board. Owner: Chief Sustainability "
          "Officer. Target: 30 June 2026."),
    ("h2", "Item 5 — Borrowing programme"),
    ("p", "The Board approved an enhancement of the working capital borrowing limit to "
          "INR 450 crore. The CFO was DIRECTED to report the debt service coverage ratio to "
          "the Board at each subsequent meeting."),
]


def build_minutes(path: Path) -> None:
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("m_h1", parent=styles["Heading1"], fontSize=13, spaceAfter=8)
    h2 = ParagraphStyle("m_h2", parent=styles["Heading2"], fontSize=11, spaceBefore=10)
    body = ParagraphStyle(
        "m_body", parent=styles["Normal"], fontSize=9.5, leading=13.5, spaceAfter=6
    )

    doc = SimpleDocTemplate(
        str(path), pagesize=A4,
        leftMargin=20 * mm, rightMargin=20 * mm, topMargin=20 * mm, bottomMargin=20 * mm,
        title=f"Board Minutes - {COMPANY}",
    )
    story = []
    for kind, text in MINUTES:
        if kind == "pb":
            story.append(PageBreak())
        elif kind == "h1":
            story.append(Paragraph(text, h1))
        elif kind == "h2":
            story.append(Paragraph(text, h2))
        else:
            story.append(Paragraph(text, body))
    story.append(Spacer(1, 12))
    doc.build(story)


# --- 2. Board deck (PPTX) -----------------------------------------------------

SLIDES = [
    (
        "119th Board Meeting — 21 August 2026",
        [f"{COMPANY}", "Board pack for the quarter ended 30 June 2026", "Strictly confidential"],
        "",
    ),
    (
        "Agenda",
        [
            "1. Confirmation of minutes of the 118th meeting",
            "2. Financial results — quarter ended 30 June 2026",
            "3. Business review",
            "4. Kalyani Auto Components — revised proposal (for approval)",
            "5. Capital expenditure — Nashik line 3 (for approval)",
            "6. Risk management update",
            "7. Any other business",
        ],
        "",
    ),
    (
        "Q1 FY27 — financial highlights",
        [
            "Revenue INR 1,182 cr, up 4.1% year on year",
            "EBITDA INR 168 cr; margin 14.2% against 17.8% in Q4 FY26",
            "Profit after tax INR 61 cr, down 22% year on year",
            "Net debt INR 1,240 cr against INR 890 cr at 31 March 2026",
            "Order book INR 3,410 cr, up 11% year on year",
        ],
        "Margin decline is attributable primarily to the one-off provision on the Nashik "
        "inventory write-down. Underlying margin excluding the provision would be 16.9%.",
    ),
    (
        "Business review — commercial",
        [
            "Volumes flat quarter on quarter; realisation down 2.4%",
            "Two large customers renegotiated pricing during the quarter",
            "Sales headcount reduced from 240 to 198 over the last two quarters",
            "New product pipeline on track for H2 FY27 launch",
        ],
        "",
    ),
    (
        "Kalyani Auto Components — revised proposal",
        [
            "Acquisition of 74% equity stake",
            "Enterprise value INR 385 cr, up from INR 340 cr in the original proposal",
            "Funding: INR 260 cr debt, INR 125 cr internal accruals",
            "Expected to be earnings accretive from FY29",
            "For the approval of the Board",
        ],
        "The independent valuation opinion has been commissioned and is expected shortly. "
        "Management recommends approval subject to receipt of the opinion.",
    ),
    (
        "Capital expenditure — Nashik line 3",
        [
            "Proposed outlay INR 210 cr over 18 months",
            "Capacity addition 40,000 tonnes per annum",
            "Indicative payback 5.5 years at current realisation",
            "For the approval of the Board",
        ],
        "Payback assumes realisation recovers to FY26 levels from H2 FY27.",
    ),
    (
        "Risk management update",
        [
            "Cyber security: phishing incident in June 2026 contained; no data exfiltration confirmed",
            "Customer concentration: top three customers now 41% of revenue, up from 34%",
            "Regulatory: show-cause notice received from the State Pollution Control Board, Nashik",
            "Working capital: receivables above 180 days at INR 78 cr",
        ],
        "",
    ),
    (
        "Debt service coverage",
        [
            "DSCR for the quarter 1.28x against covenant floor of 1.20x",
            "Headroom has narrowed from 1.61x at 31 March 2026",
        ],
        "",
    ),
]


def build_deck(path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[5]  # Title only

    for title, bullets, notes in SLIDES:
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = PPt(28)

        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.4), Inches(4.4))
        frame = box.text_frame
        frame.word_wrap = True
        for index, bullet in enumerate(bullets):
            para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = PPt(16)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(path))


# --- 3. Financial MIS (XLSX) --------------------------------------------------

PL_ROWS = [
    ["Particulars (INR crore)", "Q1 FY27", "Q4 FY26", "Q1 FY26", "QoQ %", "YoY %"],
    ["Revenue from operations", 1182, 1229, 1135, -3.8, 4.1],
    ["Other income", 14, 18, 12, -22.2, 16.7],
    ["Cost of materials consumed", 703, 704, 651, -0.1, 8.0],
    ["Employee benefits expense", 141, 138, 132, 2.2, 6.8],
    ["Nashik inventory write-down (one-off)", 32, 0, 0, 0.0, 0.0],
    ["Other expenses", 152, 186, 156, -18.3, -2.6],
    ["EBITDA", 168, 219, 208, -23.3, -19.2],
    ["EBITDA margin %", 14.2, 17.8, 18.3, 0.0, 0.0],
    ["Finance costs", 41, 33, 29, 24.2, 41.4],
    ["Depreciation", 48, 47, 44, 2.1, 9.1],
    ["Profit before tax", 79, 139, 135, -43.2, -41.5],
    ["Profit after tax", 61, 103, 78, -40.8, -21.8],
]

BS_ROWS = [
    ["Particulars (INR crore)", "30 Jun 2026", "31 Mar 2026", "Change"],
    ["Trade receivables — total", 612, 548, 64],
    ["  of which over 180 days", 78, 41, 37],
    ["Inventories", 428, 396, 32],
    ["Cash and equivalents", 96, 174, -78],
    ["Gross borrowings", 1336, 1064, 272],
    ["Net debt", 1240, 890, 350],
    ["Net worth", 1985, 1948, 37],
    ["Net debt to EBITDA (TTM)", 1.72, 1.11, 0.61],
]

COV_ROWS = [
    ["Covenant", "Floor / Cap", "Q1 FY27", "Q4 FY26", "Headroom"],
    ["Debt service coverage ratio", "min 1.20x", 1.28, 1.61, 0.08],
    ["Net debt to EBITDA", "max 2.50x", 1.72, 1.11, 0.78],
    ["Interest coverage ratio", "min 3.00x", 4.10, 6.64, 1.10],
    ["Current ratio", "min 1.10x", 1.14, 1.31, 0.04],
]

SEG_ROWS = [
    ["Segment", "Q1 FY27 revenue", "Q1 FY26 revenue", "Q1 FY27 EBIT", "Q1 FY26 EBIT"],
    ["Automotive components", 684, 622, 71, 88],
    ["Industrial castings", 341, 356, 38, 44],
    ["Aftermarket and services", 157, 157, 22, 20],
]


def build_mis(path: Path) -> None:
    workbook = openpyxl.Workbook()

    sheets = [
        ("P&L", PL_ROWS),
        ("Balance Sheet", BS_ROWS),
        ("Covenants", COV_ROWS),
        ("Segments", SEG_ROWS),
    ]

    first = workbook.active
    first.title = sheets[0][0]
    for row in sheets[0][1]:
        first.append(row)

    for name, rows in sheets[1:]:
        sheet = workbook.create_sheet(name)
        for row in rows:
            sheet.append(row)

    workbook.save(str(path))


# --- 4. Risk register + internal audit (DOCX) ---------------------------------

RISK_DOC = [
    ("h1", f"{COMPANY} — Enterprise Risk Register"),
    ("h2", "Quarter ended 30 June 2026 — presented to the Risk Management Committee"),
    ("h2", "R-01  Customer concentration"),
    ("p", "Severity: High (raised from Medium). The top three customers accounted for 41 per "
          "cent of revenue in Q1 FY27, against 34 per cent in Q1 FY26. Two of the three "
          "renegotiated pricing during the quarter. Mitigation: the new product pipeline is "
          "expected to broaden the customer base from H2 FY27. Owner: Chief Commercial Officer."),
    ("h2", "R-02  Working capital and liquidity"),
    ("p", "Severity: High. Receivables above 180 days stood at INR 78 crore at 30 June 2026 "
          "against INR 41 crore at 31 March 2026. Cash and equivalents reduced to INR 96 crore. "
          "The debt service coverage ratio was 1.28 times against a covenant floor of 1.20 times. "
          "Owner: Chief Financial Officer."),
    ("h2", "R-03  Cyber security"),
    ("p", "Severity: Medium. A phishing incident was detected on 9 June 2026 affecting four "
          "mailboxes in the finance function. The incident was contained within 36 hours. "
          "Forensic review is ongoing and no data exfiltration has been confirmed to date. "
          "The company's cyber insurance cover is INR 25 crore. Owner: Chief Information Officer."),
    ("h2", "R-04  Regulatory and environmental"),
    ("p", "Severity: High. A show-cause notice was received from the State Pollution Control "
          "Board on 2 July 2026 in respect of effluent discharge parameters at the Nashik plant. "
          "The maximum penalty under the applicable rules is INR 12 crore, and the notice "
          "contemplates the possibility of a direction to suspend operations at the affected "
          "line. A reply is due within 30 days. Owner: Head of Environment, Health and Safety."),
    ("h2", "R-05  Key person dependency"),
    ("p", "Severity: Medium. No documented succession plan exists for the Managing Director or "
          "Chief Financial Officer positions. Owner: Head of Human Resources."),
    ("pb", ""),
    ("h1", "Internal Audit — Summary of open observations"),
    ("h2", "Quarter ended 30 June 2026"),
    ("h2", "IA-2025-07  Segregation of duties — procurement, Nashik"),
    ("p", "Rating: High. Originally raised in the audit cycle ended March 2025 and repeated in "
          "the cycle ended March 2026. The same user identifiers retain the ability to create a "
          "vendor master record and to release a purchase order. Management response: "
          "remediation is dependent on the ERP authorisation redesign, which has been deferred "
          "pending the IT roadmap. Revised target date: 31 March 2027."),
    ("h2", "IA-2026-02  Inventory valuation — Nashik"),
    ("p", "Rating: High. Slow-moving inventory of INR 32 crore was identified as requiring "
          "provision. The provision has been recognised in Q1 FY27. Internal audit observes that "
          "the slow-moving condition was identifiable from the ageing report as at 31 December "
          "2025 and was not escalated at that time."),
    ("h2", "IA-2026-05  Customer credit limits"),
    ("p", "Rating: Medium. Credit limits were exceeded without documented approval in 14 of 40 "
          "instances tested, covering INR 61 crore of exposure. Management response: the credit "
          "policy is under revision."),
]


def build_risk_doc(path: Path) -> None:
    document = docx.Document()
    document.styles["Normal"].font.size = Pt(10.5)

    for kind, text in RISK_DOC:
        if kind == "pb":
            document.add_page_break()
        elif kind == "h1":
            document.add_heading(text, level=1)
        elif kind == "h2":
            document.add_heading(text, level=2)
        else:
            document.add_paragraph(text)

    document.save(str(path))


# --- entry point --------------------------------------------------------------


def main(destination: Path) -> None:
    destination.mkdir(parents=True, exist_ok=True)

    targets = [
        ("Board Minutes - 118th Meeting.pdf", build_minutes),
        ("Board Deck - 119th Meeting.pptx", build_deck),
        ("Financial MIS - Q1 FY27.xlsx", build_mis),
        ("Risk Register and Internal Audit Report.docx", build_risk_doc),
    ]

    for name, builder in targets:
        path = destination / name
        builder(path)
        print(f"  {name}  ({path.stat().st_size / 1024:.0f} KB)")

    print(f"\nSample board pack written to {destination.resolve()}")


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("sample_pack")
    main(out)

"""PPTX ingestion via python-pptx.

Board decks carry a large share of the substance, and speaker notes often hold
the caveats that the slide body omits - both are indexed. Slide numbers are
exact, so deck citations are the most precise in the pack.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import Segment, normalise


def parse(path: Path) -> list[Segment]:
    prs = Presentation(str(path))
    segments: list[Segment] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        title = _title(slide)
        body_parts: list[str] = []
        tables: list[str] = []

        for shape in _walk(slide.shapes):
            if shape.has_text_frame:
                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text and text != title:
                    body_parts.append(text)
            if getattr(shape, "has_table", False):
                tables.append(_render_table(shape.table))

        body = normalise("\n".join(body_parts))
        if title or body:
            segments.append(
                Segment(
                    page=slide_no,
                    text=normalise(f"{title}\n{body}" if title else body),
                    locator=f"Slide {slide_no}",
                    heading=title,
                    kind="body",
                )
            )

        for t_index, table in enumerate(tables, start=1):
            if table:
                segments.append(
                    Segment(
                        page=slide_no,
                        text=table,
                        locator=f"Slide {slide_no}, table {t_index}",
                        heading=title,
                        kind="table",
                    )
                )

        notes = _notes(slide)
        if notes:
            segments.append(
                Segment(
                    page=slide_no,
                    text=notes,
                    locator=f"Slide {slide_no} (speaker notes)",
                    heading=title,
                    kind="notes",
                )
            )

    return segments


def _walk(shapes):
    """Flatten grouped shapes - decks nest content in groups constantly."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk(shape.shapes)
        else:
            yield shape


def _title(slide) -> str:
    try:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except (AttributeError, ValueError):
        pass
    return ""


def _notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        return ""
    return normalise(frame.text or "")


def _render_table(table) -> str:
    lines: list[str] = []
    for row in table.rows:
        cells = [c.text.strip().replace("\n", " ") for c in row.cells]
        if any(cells):
            lines.append(" | ".join(cells))
    return "\n".join(lines)
